import os
import io
import streamlit as st
import pandas as pd
import plotly.express as px
import plotly.graph_objects as go

# ניסיון לטעון את ספריית יצירת המצגות - לא יקרוס אם חסר
try:
    from pptx import Presentation
    from pptx.util import Inches
    PPTX_AVAILABLE = True
except ImportError:
    PPTX_AVAILABLE = False

# --- הגדרות תצוגה בסיסיות ---
st.set_page_config(page_title="PE Fund Analytics", layout="wide", page_icon="📈")

# --- מנגנון טעינת נתונים חכם (היברידי) ---
@st.cache_data
def load_excel_data(file_or_path):
    try:
        xls = pd.ExcelFile(file_or_path)
    except Exception as e:
        return pd.DataFrame(), pd.DataFrame(), []
        
    sheet_names = xls.sheet_names
    combined_sheet = next((s for s in sheet_names if 'Combined' in s), None)
    cleaned_sheet = next((s for s in sheet_names if 'Cleaned' in s), None)
    
    df_combined = pd.DataFrame()
    df_cleaned = pd.DataFrame()
    
    if combined_sheet:
        df_combined = pd.read_excel(xls, sheet_name=combined_sheet)
        if 'Date of Investment' in df_combined.columns and 'Exit/Current Date' in df_combined.columns:
            df_combined['Date of Investment'] = pd.to_datetime(df_combined['Date of Investment'], errors='coerce')
            df_combined['Exit/Current Date'] = pd.to_datetime(df_combined['Exit/Current Date'], errors='coerce')
            df_combined['Hold Period (Years)'] = (df_combined['Exit/Current Date'] - df_combined['Date of Investment']).dt.days / 365.25
            df_combined['Vintage Year'] = df_combined['Date of Investment'].dt.year
            
        if 'Gross Multiple' in df_combined.columns:
            df_combined['Gross Multiple'] = pd.to_numeric(df_combined['Gross Multiple'], errors='coerce')
            
    if cleaned_sheet:
        df_cleaned = pd.read_excel(xls, sheet_name=cleaned_sheet)
        
    return df_combined, df_cleaned, sheet_names

# --- פונקציה לייצוא PPTX ---
def create_ppt_report(fund_name, deals_count, fig_waterfall):
    prs = Presentation()
    
    # שקף 1: כותרת
    title_slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(title_slide_layout)
    slide.shapes.title.text = f"Tear Sheet: {fund_name}"
    slide.placeholders[1].text = f"מבוסס על ניתוח של {deals_count} עסקאות\nהופק אוטומטית מ-PE Analytics"
    
    # שקף 2: גרף (אם יש)
    if fig_waterfall:
        blank_slide_layout = prs.slide_layouts[5]
        slide2 = prs.slides.add_slide(blank_slide_layout)
        slide2.shapes.title.text = "Value Creation Bridge (Gross MOIC)"
        try:
            img_bytes = io.BytesIO()
            fig_waterfall.write_image(img_bytes, format='png', width=800, height=450)
            img_bytes.seek(0)
            slide2.shapes.add_picture(img_bytes, Inches(1), Inches(2), width=Inches(8))
        except Exception as e:
            slide2.shapes.title.text = "שגיאה בטעינת הגרף למצגת (יש לוודא התקנת kaleido)"
    
    ppt_io = io.BytesIO()
    prs.save(ppt_io)
    ppt_io.seek(0)
    return ppt_io

# ==========================================
# בניית ה-UI
# ==========================================
st.sidebar.title("ניווט והגדרות")

# בדיקה האם הקובץ קיים מקומית
FILE_PATH = 'Fund Proprietary Database.xlsx'
file_to_load = None

if os.path.exists(FILE_PATH):
    file_to_load = FILE_PATH
    st.sidebar.success("✅ קובץ נתונים מקומי זוהה ונטען אוטומטית.")
else:
    st.sidebar.warning("⚠️ לא זוהה קובץ מקומי.")
    file_to_load = st.sidebar.file_uploader("העלה את קובץ האקסל כאן:", type=["xlsx"])

# רק אם יש קובץ לעבוד איתו (מקומי או שהועלה הרגע)
if file_to_load is not None:
    df, df_clean, all_sheets = load_excel_data(file_to_load)
    
    if df.empty:
        st.error("הקובץ נטען אך הגיליונות הרלוונטיים (Combined/Cleaned) חסרים בו.")
    else:
        page = st.sidebar.radio("בחר מסך ניתוח:", 
                                ["מבט על (Macro Overview)", 
                                 "מגמות ומפות חום (Trends & Maps)",
                                 "פרופיל קרן (GP Tear Sheet) 📊", 
                                 "מגרש המשחקים (Deal Explorer)",
                                 "תשאול חופשי (AI Chat) 🤖"])
        
        st.sidebar.markdown("---")

        # --- מסך 1: מבט על ---
        if page == "מבט על (Macro Overview)":
            st.title("מבט על - ביצועי תעשיית ה-PE")
            col1, col2, col3, col4 = st.columns(4)
            col1.metric("סה\"כ עסקאות", f"{len(df):,}")
            if 'Gross Multiple' in df.columns: col2.metric("מכפיל (MOIC) ממוצע", f"{df['Gross Multiple'].mean():.2f}x")
            if 'Fund Name' in df.columns: col3.metric("מספר קרנות", f"{df['Fund Name'].nunique()}")
            if 'Hold Period (Years)' in df.columns: col4.metric("זמן החזקה ממוצע", f"{df['Hold Period (Years)'].mean():.1f} שנים")
            
            st.markdown("---")
            if 'Gross Multiple' in df.columns and 'Fund Name' in df.columns:
                st.subheader("דירוג מנהלים: תשואה לעומת יחס הפסד")
                fund_stats = df.dropna(subset=['Gross Multiple', 'Fund Name']).groupby('Fund Name').agg(
                    Avg_MOIC=('Gross Multiple', 'mean'), Deal_Count=('Gross Multiple', 'count'),
                    Loss_Ratio=('Gross Multiple', lambda x: (x < 1.0).mean() * 100)
                ).reset_index()
                fund_stats = fund_stats[fund_stats['Deal_Count'] > 3] 
                fig_scatter = px.scatter(fund_stats, x="Loss_Ratio", y="Avg_MOIC", size="Deal_Count", color="Fund Name", hover_name="Fund Name", title="גודל הבועה = מספר עסקאות")
                fig_scatter.add_hline(y=1.0, line_dash="dot", line_color="red")
                st.plotly_chart(fig_scatter, use_container_width=True)

        # --- מסך 2: מגמות ומפות חום ---
        elif page == "מגמות ומפות חום (Trends & Maps)":
            st.title("מגמות מאקרו-כלכליות וגיאוגרפיה")
            tab1, tab2 = st.tabs(["שנות בציר וסקטורים", "מפת עולם"])
            
            with tab1:
                st.subheader("החזר ממוצע (MOIC) לפי שנת השקעה וסקטור")
                if 'Vintage Year' in df.columns and 'Industry' in df.columns:
                    valid_years = df[(df['Vintage Year'] >= 1990) & (df['Vintage Year'] <= 2024)].copy()
                    heatmap_data = valid_years.groupby(['Vintage Year', 'Industry'])['Gross Multiple'].mean().reset_index()
                    top_industries = valid_years['Industry'].value_counts().nlargest(10).index
                    heatmap_data = heatmap_data[heatmap_data['Industry'].isin(top_industries)]
                    fig_heat = px.density_heatmap(heatmap_data, x="Vintage Year", y="Industry", z="Gross Multiple", histfunc="avg", color_continuous_scale="Viridis")
                    st.plotly_chart(fig_heat, use_container_width=True)
            
            with tab2:
                st.subheader("מפת ביצועים גלובלית")
                if 'Country' in df.columns:
                    map_data = df.groupby('Country').agg(Avg_MOIC=('Gross Multiple', 'mean'), Deal_Count=('Gross Multiple', 'count')).reset_index()
                    fig_map = px.choropleth(map_data, locations="Country", locationmode="country names", color="Avg_MOIC", hover_name="Country", color_continuous_scale="Blues")
                    st.plotly_chart(fig_map, use_container_width=True)

        # --- מסך 3: פרופיל קרן (וייצוא PPTX) ---
        elif page == "פרופיל קרן (GP Tear Sheet) 📊":
            st.title("פרופיל מנהל קרן ויצירת ערך")
            if not df_clean.empty and 'Fund Name' in df_clean.columns:
                funds_list = sorted(df_clean['Fund Name'].dropna().unique())
                selected_fund = st.selectbox("בחר קרן:", funds_list)
                fund_data = df_clean[df_clean['Fund Name'] == selected_fund]
                
                fig_waterfall = None
                vc_cols = ['Revenue_MOIC', 'Margin_MOIC', 'Multiple_MOIC', 'Net_Debt_MOIC']
                
                if all(c in fund_data.columns for c in vc_cols):
                    fund_vc_data = fund_data.dropna(subset=vc_cols)
                    if len(fund_vc_data) > 0:
                        st.subheader(f"גשר יצירת ערך - ממוצע לעסקה בקרן")
                        avg_entry = 1.0
                        avg_rev = fund_vc_data['Revenue_MOIC'].mean()
                        avg_mar = fund_vc_data['Margin_MOIC'].mean()
                        avg_mul = fund_vc_data['Multiple_MOIC'].mean()
                        avg_debt = fund_vc_data['Net_Debt_MOIC'].mean()
                        avg_total = avg_entry + avg_rev + avg_mar + avg_mul + avg_debt
                        
                        fig_waterfall = go.Figure(go.Waterfall(
                            orientation="v", measure=["absolute", "relative", "relative", "relative", "relative", "total"],
                            x=["Entry", "Revenue Growth", "Margin Expansion", "Multiple Arbitrage", "Debt/Cash Flow", "Current/Exit"],
                            text=[f"{avg_entry:.2f}x", f"{avg_rev:+.2f}x", f"{avg_mar:+.2f}x", f"{avg_mul:+.2f}x", f"{avg_debt:+.2f}x", f"{avg_total:.2f}x"],
                            y=[avg_entry, avg_rev, avg_mar, avg_mul, avg_debt, avg_total],
                            connector={"line":{"color":"rgb(63, 63, 63)"}}
                        ))
                        st.plotly_chart(fig_waterfall, use_container_width=True)
                
                st.markdown("---")
                if PPTX_AVAILABLE:
                    st.write("📥 **הפקת דוח להנהלה:** לחץ ליצירת מצגת PowerPoint אוטומטית עם ביצועי הקרן.")
                    if st.button("הכן מצגת PPTX"):
                        with st.spinner("מכין מצגת..."):
                            ppt_file = create_ppt_report(selected_fund, len(fund_data), fig_waterfall)
                            st.download_button("⬇️ הורד קובץ", data=ppt_file, file_name=f"{selected_fund}_Report.pptx", mime="application/vnd.openxmlformats-officedocument.presentationml.presentation")
                else:
                    st.error("ספריית python-pptx לא מותקנת ולכן לא ניתן לייצר מצגת. הרץ pip install python-pptx.")
                    
                st.subheader("עסקאות הקרן")
                cols_to_show = [c for c in ['Portfolio Company', 'Industry', 'Gross Multiple'] if c in df.columns]
                st.dataframe(df[df['Fund Name'] == selected_fund][cols_to_show].sort_values(by='Gross Multiple', ascending=False), use_container_width=True)

        # --- מסך 4: מגרש המשחקים ---
        elif page == "מגרש המשחקים (Deal Explorer)":
            st.title("מגרש המשחקים - חקר עסקאות")
            selected_sector = st.sidebar.multiselect("סנן לפי סקטור:", options=df['Industry'].dropna().unique())
            max_val = float(df['Gross Multiple'].max()) if not df['Gross Multiple'].isna().all() else 15.0
            min_moic, max_moic = st.sidebar.slider("טווח מכפיל:", 0.0, min(max_val, 20.0), (0.0, 10.0))
            
            filtered_df = df.copy()
            if selected_sector: filtered_df = filtered_df[filtered_df['Industry'].isin(selected_sector)]
            filtered_df = filtered_df[(filtered_df['Gross Multiple'] >= min_moic) & (filtered_df['Gross Multiple'] <= max_moic)]
            
            st.write(f"מציג {len(filtered_df)} עסקאות שעונות לתנאים.")
            fig_deals = px.scatter(filtered_df, x="Hold Period (Years)", y="Gross Multiple", hover_name="Portfolio Company", color="Industry")
            st.plotly_chart(fig_deals, use_container_width=True)

        # --- מסך 5: תשאול AI חופשי ---
        elif page == "תשאול חופשי (AI Chat) 🤖":
            st.title("🤖 שיחה עם מסד הנתונים")
            if "messages" not in st.session_state:
                st.session_state.messages = []

            for message in st.session_state.messages:
                with st.chat_message(message["role"]):
                    st.markdown(message["content"])

            if prompt := st.chat_input("לדוגמה: הראה לי עסקאות מפסידות של KKR"):
                st.session_state.messages.append({"role": "user", "content": prompt})
                with st.chat_message("user"): st.markdown(prompt)

                with st.chat_message("assistant"):
                    if "הפסד" in prompt or "kkr" in prompt.lower():
                        response = "להלן העסקאות המפסידות של KKR במסד הנתונים:"
                        st.markdown(response)
                        losses = df[(df['Fund Name'].str.contains('KKR', na=False)) & (df['Gross Multiple'] < 1.0)]
                        st.dataframe(losses[['Portfolio Company', 'Fund Name', 'Gross Multiple']])
                    elif "תוכנה" in prompt or "software" in prompt.lower():
                        response = "קרנות עם הביצועים הטובים ביותר בסקטור התוכנה:"
                        st.markdown(response)
                        sw_df = df[df['Industry'].str.contains('Software', case=False, na=False)]
                        top_sw = sw_df.groupby('Fund Name')['Gross Multiple'].mean().nlargest(5).reset_index()
                        st.table(top_sw)
                    else:
                        response = f"אני מנוע הדגמה (Mock). זיהיתי ששאלת: '{prompt}'. כדי לקבל תשובות על כל שאלה, יש לחבר מפתח OpenAI למערכת."
                        st.markdown(response)
                st.session_state.messages.append({"role": "assistant", "content": response})

else:
    st.info("אנא העלה את קובץ האקסל (Fund Proprietary Database.xlsx) או ודא שהוא שמור באותה תיקייה של הקוד.")
